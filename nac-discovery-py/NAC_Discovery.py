import json,os  
import boto3
from datetime import *
import json, logging
import pprint,re
from elasticsearch import Elasticsearch, helpers
from opensearchpy import OpenSearch,helpers, RequestsHttpConnection
import requests
from requests_aws4auth import AWS4Auth
import urllib.parse
from botocore.exceptions import ClientError
from datetime import *
import shlex,subprocess,json
from urllib.parse import unquote_plus
import elasticsearch
# import PyPDF2
from io import BytesIO
import io  
from pptx import Presentation
import fitz
from requests.auth import HTTPBasicAuth
from docx import Document
import pandas as pd
import logging
import ssl
import urllib3


logging.getLogger().setLevel(logging.INFO)
logging.info(f'date={date}')
cfn = boto3.resource('cloudformation')
def lambda_handler(event, context): 
    logging.info('lambda_handler starts...')
    print('context.invoked_function_arn',context.invoked_function_arn)
    logging.info("Lambda function ARN:".format(context.invoked_function_arn))
    runtime_region = os.environ['AWS_REGION'] 
    context_arn=context.invoked_function_arn
    u_id=context_arn.split('-')[-1]
    logging.info('u_id'.format(u_id))
    logging.info('***********************************************')
    s3 = boto3.client('s3')        
    data={}
    doc_list=[]
    check=0
    secret_data_internal = get_secret(
        'nasuni-labs-internal-'+u_id, runtime_region)
    secret_nct_nce_admin = get_secret('nasuni-labs-os-admin',runtime_region) 
    
    role = secret_data_internal['discovery_lambda_role_arn']
    username=secret_nct_nce_admin['nac_es_admin_user']
    role_data = '{"backend_roles":["' +role + '"],"hosts": [],"users": ["'+username+'"]}'
    logging.info('role_data'.format(role_data))
    with open("/tmp/"+"/data.json", "w") as write_file:
        write_file.write(role_data)
        
    link=secret_nct_nce_admin['nac_kibana_url']
    link=link[:link.index('_')]

    password=secret_nct_nce_admin['nac_es_admin_password']
    data_file_obj = '/tmp/data.json'
    merge_link = '\"https://'+link+'_opendistro/_security/api/rolesmapping/all_access\"'
    url = 'https://' + link + '_opendistro/_security/api/rolesmapping/all_access/'

    headers = {'content-type': 'application/json'}
    response = requests.put(url, auth=HTTPBasicAuth(username, password), headers=headers, data=role_data)
    logging.info("response.text {}".format(response.text))
    # share_data=call_nmc_apis(runtime_region,secret_data_internal) 
    #traversing thru each file metadata and data and adding those fields into dictionary
    
    #traversing thru each file metadata and data and adding those fields into dictionary
    bucket_name='nasuni-share-data-bucket-storage'    
    # s3 = boto3.client('s3')
    print(bucket_name)
    
    # List all of the files in the S3 bucket
    response = s3.list_objects(Bucket=bucket_name)
    
    # Read the contents of each file in the S3 bucket
    print('response',response)
    bucket_folder_name=None
    share_data={}
    for obj in response['Contents']: 
        # Get the object key (i.e. the file name)
        key = obj['Key']
        bucket_folder_name=key
        # print('bucket_folder_name',bucket_folder_name)
        print('key',key)  
        if u_id in key:
            print('found',key)
            nmc_api_filename=os.path.basename(key)
            print('nmc_api_filename',nmc_api_filename)
            s3.download_file(bucket_name, key, '/tmp/'+nmc_api_filename)
            
            with open('/tmp/'+nmc_api_filename, 'r') as f2:
                if 'nmc_api_data_v_share_name' in '/tmp/'+nmc_api_filename:
                    share_data['name'] = f2.read().split(',')
                else:
                    share_data['path'] = f2.read().split(',')
            
                # print(data_file)
            logging.info('deleting folder from s3 bucket nasuni-share-data-bucket-storage')
            # s3.delete_object(Bucket=bucket_name, Key=key)
    print(share_data)
    logging.info(share_data)
    
    for record in event['Records']:
        logging.info(record)
        data['dest_bucket'] = record['s3']['bucket']['name']
        data['object_key'] = unquote_plus(record['s3']['object']['key'])
        data['size'] = str(record['s3']['object'].get('size', -1))
        file_name=os.path.basename(data['object_key'])
        data['file_name'] = file_name
        data['event_name'] = record['eventName']
        data['event_time'] = record['eventTime']
        data['awsRegion'] = record['awsRegion']

        data['extension'] = file_name[file_name.index('.')+1:]
        data['volume_name'] = secret_data_internal['volume_name']
        
        #data['root_handle'] = secret_data_internal['root_handle'].replace('.','_').lower()
        data['root_handle'] = re.sub('[!@#$%^&*()+?=,<>/.]', '-', secret_data_internal['root_handle']).lower()
        data['source_bucket'] = secret_data_internal['discovery_source_bucket']
        logging.info("data['object_key'] = {}".format(data['object_key']))  
        logging.info("data['dest_bucket'] = {}".format(data['dest_bucket']))  
        obj1 = s3.get_object(Bucket=data['dest_bucket'], Key=data['object_key'])
        if  data['extension'] in ['txt','csv','docx','doc','pdf','xlsx','xls','pptx','ppt']:
                
            if data['extension'] in ['csv','txt']:
                data['content'] = obj1['Body'].read().decode('utf-8')
            elif data['extension'] == 'pdf':
                file_content = obj1['Body'].read()
                text = ""
                with fitz.open(stream=file_content, filetype="pdf") as doc:
                    
                    # iterating through pdf file pages
                    for page in range(doc.page_count):
                        # fetching & appending text to text variable of each page
                        # text += page.getText()
                        text += doc.get_page_text(page) 
                    
                data['content'] = text
            elif data['extension'] in ['docx','doc']:
               fs = obj1['Body'].read()
               sentence = str(parseDocx(fs))
               logging.info('docx data {} '.format(sentence))
               data['content'] = sentence
            elif data['extension'] in ['xlsx','xls']:
                file_content = obj1['Body'].read()
                read_excel_data = io.BytesIO(file_content)
                df = pd.read_excel(read_excel_data) 
                df = df.to_string() 
                logging.info('xlsx data {}'.format(df))
                data['content'] = df 
            elif data['extension'] in ['pptx','ppt']:
                print('data[extension] elif',data['extension'])
                pptx_content = obj1['Body'].read()
                ppt = Presentation(io.BytesIO(pptx_content))
                pptx_data=''
                for slide in ppt.slides:
                    for shape in slide.shapes:
                        if not shape.has_text_frame:
                            continue
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                pptx_data+=run.text
                logging.info("pptx data {}".format(pptx_data))
                data['content'] = pptx_data
        else:
            data['content'] =data['file_name']
        share_path_last_element=None
        list_after_index=None
        if share_data != None:
            if share_data['name']  and share_data['path']:
                for name,path in zip(share_data['name'],share_data['path']):
                    
                    if path in data['object_key']:
                        share_path_last_element=path.split('/')[-1] 
                        logging.info('148 share_path_last_element {}'.format(share_path_last_element))
                        full_path=data['object_key']

                        full_path_with_share_name=full_path.replace(path,'/'+name)
                        logging.info('full_path_with_share_name {}'.format(full_path_with_share_name))
                        index_of_last_element=full_path_with_share_name.index(name)

                        list_after_index=full_path_with_share_name[index_of_last_element:]
                    
        if secret_data_internal['web_access_appliance_address']!='not_found':
            if share_path_last_element != None:
                if share_data['name'] and share_data['path'] and share_path_last_element in data['object_key']:
                    # data['access_url']='https://'+secret_data_internal['web_access_appliance_address']+'/fs/view/'+secret_data_internal['share_name']+'/'+list_after_index
                    data['access_url']='https://'+secret_data_internal['web_access_appliance_address']+'/fs/view/'+list_after_index
            else:
                data['access_url']='https://'+secret_data_internal['web_access_appliance_address']+'/fs/view/'+data['volume_name']+'/'+'/'.join(data['object_key'].split('/')[3:])
        else:
            data['access_url']=secret_data_internal['web_access_appliance_address']
        
        logging.info('access_url = {}'.format(data['access_url']))
        percent_20_url=data['access_url'].replace(' ','%20')
        logging.info('After appending percent 20 url = {}'.format(percent_20_url))
        data['access_url']=percent_20_url
        logging.info('secret_data_internal = {}'.format(secret_data_internal))
        es_obj = launch_es(secret_nct_nce_admin['nac_es_url'],data['awsRegion'])
        
        check=connect_es(es_obj,data['root_handle'], data) 
    #Deletion of folder from s3
    if check == 0:
        logging.info('Insertion into ES success.Hence deleting s3 bucket folder')
        del_s3_folder(data['object_key'],data['dest_bucket'])
    else:
        logging.info('Not deleting the s3 bucket folder all data not got loaded into ES.') 
    logging.info('Deleting the files from /tmp/ folder')
    subprocess.call('rm -rf /tmp/*', shell=True)
    
    
    logging.info('lambda_handler ends...')
    
def call_nmc_apis(region,internal_secret):
    user_secret = get_secret(internal_secret['user_secret_name'],region)
    logging.info('user_secret = {}'.format(user_secret))
    endpoint=user_secret['nmc_api_endpoint']
    username=user_secret['nmc_api_username']
    password=user_secret['nmc_api_password']
    if not os.environ.get('PYTHONHTTPSVERIFY', '') and getattr(ssl, '_create_unverified_context', None):
        ssl._create_default_https_context = ssl._create_unverified_context
    url = 'https://' + endpoint + '/api/v1.1/auth/login/'
    logging.info(url)
    try :
        values = {'username': username, 'password': password}
        data = urllib.parse.urlencode(values).encode("utf-8")
        logging.info(data)
        response = urllib.request.urlopen(url, data, timeout=5)
        logging.info(response)
        result = json.loads(response.read().decode('utf-8'))
        logging.info(result)
    except Exception as e:
        logging.error('ERROR: {0}'.format(str(e)))
        return
        
    urllib3.disable_warnings()
    headers = {
        'Accept': 'application/json',
        'Authorization': 'Token {}'.format(result['token'])
    }
    print(headers)
    vv_guid=None
    try:
        r = requests.get('https://' + endpoint + '/api/v1.1/volumes/', headers = headers,verify=False)
    except requests.exceptions.RequestException as err:
        logging.error ("OOps: Something Else {}".format(err))
    except requests.exceptions.HTTPError as errh:
        logging.error ("Http Error: {}".format(errh))
    except requests.exceptions.ConnectionError as errc:
        logging.error ("Error Connecting: {}".format(errc))
    except requests.exceptions.Timeout as errt:
        logging.error ("Timeout Error: {}".format(errt))
    except Exception as e:
        logging.error('ERROR: {0}'.format(str(e)))
    
    # print(r.json())
    for i in r.json()['items']:
        if i['name'] == internal_secret['volume_name']:
            vv_guid = i['guid']
            print(vv_guid)
    try:
        r = requests.get('https://' + endpoint + '/api/v1.1/volumes/filers/shares/', headers = headers,verify=False)
    except requests.exceptions.RequestException as err:
        logging.error ("OOps: Something Else {}".format(err))
    except requests.exceptions.HTTPError as errh:
        logging.error ("Http Error: {}".format(errh))
    except requests.exceptions.ConnectionError as errc:
        logging.error ("Error Connecting: {}".format(errc))
    except requests.exceptions.Timeout as errt:
        logging.error ("Timeout Error: {}".format(errt))
    except Exception as e:
        logging.error('ERROR: {0}'.format(str(e)))
    
    share_data={}
    name=[]
    path=[]
    for i in r.json()['items']:
        if i['volume_guid'] == vv_guid and i['path']!='\\' and i['browser_access']==True:
            name.append(r""+i['name'].replace('\\','/'))
            path.append(r""+i['path'].replace('\\','/'))
            

    share_data['name']=name
    share_data['path']=path
    logging.info(share_data) 
    return share_data
    
def parseDocx(data):
    data = io.BytesIO(data)
    document = Document(docx = data)
    content = ''
    for para in document.paragraphs:
        data = para.text
        content+= data
    return content

def del_s3_folder(full_path,dest_bucket):
    logging.info("Full Path:-".format(full_path))
    path=os.path.dirname(full_path)
    logging.info("Folder Path:-".format(path))
    s3 = boto3.resource('s3') 
    bucket = s3.Bucket(dest_bucket)
    bucket.objects.filter(Prefix=path).delete()
    

def launch_es(es_url,region):

    service = 'es'
    credentials = boto3.Session().get_credentials()
    awsauth = AWS4Auth(credentials.access_key, credentials.secret_key, region, service, session_token=credentials.token)
    # es = Elasticsearch(hosts=[{'host': es_url, 'port': 443}], http_auth=awsauth, use_ssl=True, verify_certs=True)
    # es = Elasticsearch(hosts=[{'host': es_url, 'port': 443}], http_auth=awsauth, verify_certs=True)
    es = OpenSearch(hosts=[{'host': es_url, 'port': 443}], http_auth=awsauth, use_ssl=True, verify_certs=True,connection_class = RequestsHttpConnection)
    
    return es
    
def connect_es(es,index, data):
    #CTPROJECT-125
    try:
        flag = 0
        for elem in es.cat.indices(format="json"):
            query = {"query": {"match_all": {}}}
            resp = es.search(index=elem['index'], body=query)
            for i in resp['hits']['hits']:
                idx_content = i['_source'].get('content', 0)
                idx_object_key = i['_source'].get('object_key', 0)
                if idx_content == data['content'] and idx_object_key == data['object_key']:
                    flag = 1
                    print("Indexing is doing when the idx_content and idx_object_key has matched", resp)
                    # es.index(index=i['_index'], doc_type="_doc", id=i['_id'], body=data)
                    es.index(index=i['_index'], id=i['_id'], body=data)
                    break

        if flag == 0:
            doc_list = []
            doc_list += [data]
            logging.info("\nAttempting to index the list of docs using helpers.bulk()")
            # use the helpers library's Bulk API to index list of Elasticsearch docs
            # resp = helpers.bulk(es, doc_list, index=data['root_handle'], doc_type="_doc")
            resp = helpers.bulk(es, doc_list, index=data['root_handle'])
            # print the response returned by Elasticsearch
            logging.info("helpers.bulk() RESPONSE: {}".format(resp))
            logging.info("helpers.bulk() RESPONSE: {}".format(json.dumps(resp, indent=4)))
        return 0
    except Exception as e:
        logging.error('ERROR: {0}'.format(str(e)))
        logging.error('ERROR: Unable to index line:"{0}"'.format(str(data['object_key'])))
        print(e)
        return 1

        
def get_secret(secret_name,region_name):

    secret = ''
    session = boto3.session.Session()
    client = session.client(
        service_name='secretsmanager',
        region_name=region_name,
    )

    try:
        get_secret_value_response = client.get_secret_value(
            SecretId=secret_name
        )

    except ClientError as e:
        if e.response['Error']['Code'] == 'ResourceNotFoundException':
            print("The requested secret " + secret_name + " was not found")
        elif e.response['Error']['Code'] == 'InvalidRequestException':
            print("The request was invalid due to:", e)
        elif e.response['Error']['Code'] == 'InvalidParameterException':
            print("The request had invalid params:", e)
        elif e.response['Error']['Code'] == 'DecryptionFailure':
            print("The requested secret can't be decrypted using the provided KMS key:", e)
        elif e.response['Error']['Code'] == 'InternalServiceError':
            print("An error occurred on service side:", e)
    else:
        # Secrets Manager decrypts the secret value using the associated KMS CMK
        # Depending on whether the secret was a string or binary, only one of these fields will be populated
        if 'SecretString' in get_secret_value_response:
            secret = get_secret_value_response['SecretString']

        else:
            secret = base64.b64decode(get_secret_value_response['SecretBinary'])

    return json.loads(secret)
